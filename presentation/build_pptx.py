"""
Build MBDL seminar PPTX from slides content.
Run: python3 build_pptx.py
Output: MBDL_Seminar_Final.pptx
"""

import os
from io import BytesIO

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as patches
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree

# ── Colour palette ────────────────────────────────────────────────────────────
BGU_BLUE   = RGBColor(0x00, 0x3F, 0x87)
BGU_MID    = RGBColor(0x00, 0x55, 0xB3)
ACCENT     = RGBColor(0x00, 0xA3, 0xE0)
GOLD       = RGBColor(0xF5, 0xA6, 0x23)
BG         = RGBColor(0xF8, 0xFA, 0xFD)
BG_DARK    = RGBColor(0x0B, 0x1E, 0x3D)
TEXT       = RGBColor(0x1A, 0x1A, 0x2E)
TEXT_LIGHT = RGBColor(0x4A, 0x55, 0x68)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BORDER     = RGBColor(0xD1, 0xDC, 0xF0)
ROW_ALT    = RGBColor(0xEE, 0xF4, 0xFB)
ROW_HL     = RGBColor(0xD4, 0xED, 0xFF)
ROW_HL2    = RGBColor(0xE8, 0xF5, 0xFF)

ASSETS = os.path.join(os.path.dirname(__file__), 'assets', 'figures')

W = Inches(13.333)
H = Inches(7.5)

# ── Equation rendering ────────────────────────────────────────────────────────
import re as _re

def _fix_mathtext(s):
    """Convert LaTeX commands unsupported by matplotlib mathtext to supported equivalents."""
    s = s.replace(r'\tfrac', r'\frac')
    s = s.replace(r'\dfrac', r'\frac')
    s = _re.sub(r'\\text\{', r'\\mathrm{', s)
    s = s.replace(r'\!\left', r'\left')
    s = s.replace(r'\!\right', r'\right')
    s = s.replace(r'\!\bigl', r'(')
    s = s.replace(r'\!\bigr', r')')
    s = s.replace(r'\bigl(', r'(')
    s = s.replace(r'\bigr)', r')')
    s = s.replace(r'\bigl[', r'[')
    s = s.replace(r'\bigr]', r']')
    s = s.replace(r'\bigl\{', r'\{')
    s = s.replace(r'\bigr\}', r'\}')
    s = s.replace(r'\bigl', r'')
    s = s.replace(r'\bigr', r'')
    s = s.replace(r'\;', r'\ ')
    s = s.replace(r'\!', r'')
    s = s.replace(r'\arg\min', r'\mathrm{arg\,min}')
    s = s.replace(r'\arg\max', r'\mathrm{arg\,max}')
    s = s.replace(r'\arg', r'\mathrm{arg}')
    return s


def render_eq(latex, fontsize=22, color='#003F87', bg=None, pad=0.12):
    """Render a LaTeX string to a PNG BytesIO via matplotlib mathtext."""
    latex = _fix_mathtext(latex)
    fig = plt.figure(figsize=(0.01, 0.01))
    fig.patch.set_alpha(0)
    text = fig.text(0, 0, f'${latex}$', fontsize=fontsize, color=color,
                    usetex=False, math_fontfamily='dejavusans')
    if bg:
        fig.patch.set_facecolor(bg)

    # auto-size
    fig.canvas.draw()
    bb = text.get_window_extent(renderer=fig.canvas.get_renderer())
    dpi = fig.dpi
    w = bb.width / dpi + 2 * pad
    h_in = bb.height / dpi + 2 * pad
    fig.set_size_inches(w, h_in)
    text.set_position((pad / w, pad / h_in))

    buf = BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                transparent=(bg is None), facecolor=fig.get_facecolor())
    plt.close(fig)
    buf.seek(0)
    return buf


# ── Slide helpers ─────────────────────────────────────────────────────────────
def _solid_fill(shape, color: RGBColor):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color


def _rect(slide, x, y, w, h, color: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    _solid_fill(shape, color)
    shape.line.fill.background()
    return shape


def _text_box(slide, x, y, w, h, text, size=18, bold=False, color=TEXT,
              align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.text_frame.word_wrap = wrap
    p = txb.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = 'Calibri'
    return txb


def _add_slide(prs, layout_idx=6):
    layout = prs.slide_layouts[layout_idx]  # blank
    return prs.slides.add_slide(layout)


def _top_bar(slide):
    """Gradient-effect top accent bar (solid approximation)."""
    _rect(slide, 0, 0, W, Inches(0.08), BGU_BLUE)


def _bottom_bar(slide):
    _rect(slide, 0, H - Inches(0.045), W, Inches(0.045), ACCENT)


def _slide_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_title(slide, title, y=Inches(0.32), color=BGU_BLUE, size=32):
    txb = slide.shapes.add_textbox(Inches(0.55), y, W - Inches(1.1), Inches(0.65))
    txb.text_frame.word_wrap = False
    p = txb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = color
    run.font.name = 'Calibri'
    # underline rule
    _rect(slide, Inches(0.55), y + Inches(0.62), W - Inches(1.1), Inches(0.025), BORDER)
    return txb


def _page_number(slide, n, total=16):
    _text_box(slide, W - Inches(0.9), H - Inches(0.38), Inches(0.7), Inches(0.3),
              f'{n} / {total}', size=12, color=TEXT_LIGHT, align=PP_ALIGN.RIGHT)


def _add_image(slide, path, x, y, w, h=None):
    if not os.path.exists(path):
        return
    if h:
        slide.shapes.add_picture(path, x, y, w, h)
    else:
        slide.shapes.add_picture(path, x, y, width=w)


def _add_eq_image(slide, buf, x, y, w):
    slide.shapes.add_picture(buf, x, y, width=w)


# ── Table helper ──────────────────────────────────────────────────────────────
def _add_table(slide, x, y, w, rows_data, col_widths,
               header_color=BGU_BLUE, header_text_color=WHITE,
               highlight_rows=None, highlight2_rows=None, row_size=14):
    """rows_data: list of lists of strings. First row = header."""
    cols = len(rows_data[0])
    rows = len(rows_data)
    h_row = Inches(0.33)
    total_h = h_row * rows

    table = slide.shapes.add_table(rows, cols, x, y, w, total_h).table
    table.first_row = True

    for ci, cw in enumerate(col_widths):
        table.columns[ci].width = cw

    for ri, row in enumerate(rows_data):
        for ci, cell_text in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = cell_text
            cell.text_frame.paragraphs[0].font.size = Pt(row_size)
            cell.text_frame.paragraphs[0].font.name = 'Calibri'

            # row fill
            if ri == 0:
                fill = cell.fill
                fill.solid()
                fill.fore_color.rgb = header_color
                cell.text_frame.paragraphs[0].font.color.rgb = header_text_color
                cell.text_frame.paragraphs[0].font.bold = True
            elif highlight_rows and ri in highlight_rows:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ROW_HL
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.color.rgb = BGU_BLUE
            elif highlight2_rows and ri in highlight2_rows:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ROW_HL2
                cell.text_frame.paragraphs[0].font.bold = True
            elif ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ROW_ALT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE

    return table


# ── Card helper ───────────────────────────────────────────────────────────────
def _card(slide, x, y, w, h, title, body_lines, title_color=BGU_BLUE,
          bg=WHITE, border_color=BORDER, title_size=15, body_size=13,
          left_accent=None):
    # background
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.color.rgb = border_color
    shape.line.width = Pt(0.75)

    if left_accent:
        _rect(slide, x, y, Inches(0.055), h, left_accent)

    ox = x + Inches(0.15) + (Inches(0.06) if left_accent else 0)
    ow = w - Inches(0.3) - (Inches(0.06) if left_accent else 0)

    txb = slide.shapes.add_textbox(ox, y + Inches(0.12), ow, h - Inches(0.12))
    tf = txb.text_frame
    tf.word_wrap = True

    if title:
        p = tf.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(title_size)
        p.font.color.rgb = title_color
        p.font.name = 'Calibri'
        p.text = title

    for line in body_lines:
        p = tf.add_paragraph()
        p.font.size = Pt(body_size)
        p.font.color.rgb = TEXT
        p.font.name = 'Calibri'
        p.text = line
        p.space_before = Pt(2)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide builders
# ═══════════════════════════════════════════════════════════════════════════════

def slide_01_title(prs):
    slide = _add_slide(prs)
    _slide_bg(slide, BG_DARK)

    # bottom gold bar
    _rect(slide, 0, H - Inches(0.06), W, Inches(0.06), GOLD)

    # Title
    txb = slide.shapes.add_textbox(Inches(0.9), Inches(1.5), W - Inches(1.8), Inches(1.4))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = 'Sparse Recovery and Compressed Sensing'
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = 'Calibri'

    # Subtitle
    txb2 = slide.shapes.add_textbox(Inches(0.9), Inches(3.1), W - Inches(1.8), Inches(0.7))
    tf2 = txb2.text_frame
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = 'A Model-Based Deep Learning Approach'
    run2.font.size = Pt(26)
    run2.font.color.rgb = ACCENT
    run2.font.name = 'Calibri'

    # Divider line
    _rect(slide, Inches(0.9), Inches(3.9), W - Inches(1.8), Inches(0.02),
          RGBColor(0xFF, 0xFF, 0xFF))

    # Core question box
    qbox = slide.shapes.add_shape(1, Inches(0.9), Inches(4.1), W - Inches(1.8), Inches(0.85))
    qbox.fill.solid()
    qbox.fill.fore_color.rgb = RGBColor(0x0D, 0x2D, 0x5E)
    qbox.line.fill.background()
    txb3 = slide.shapes.add_textbox(Inches(1.0), Inches(4.15), W - Inches(2.0), Inches(0.75))
    tf3 = txb3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    run3 = p3.add_run()
    run3.text = 'Core question:  How far can model-based structure take us — and when do we need data-driven flexibility?'
    run3.font.size = Pt(17)
    run3.font.color.rgb = ACCENT
    run3.font.italic = True
    run3.font.name = 'Calibri'

    # Footer
    _text_box(slide, Inches(0.9), Inches(5.25), W - Inches(1.8), Inches(0.45),
              'Etay Baron   ·   Etai Wigman',
              size=19, bold=True, color=WHITE)
    _text_box(slide, Inches(0.9), Inches(5.75), W - Inches(1.8), Inches(0.4),
              'Model-Based Deep Learning (361.2.2320)  ·  Ben-Gurion University  ·  Spring 2025–2026',
              size=15, color=RGBColor(0xAA, 0xBB, 0xCC))
    _text_box(slide, Inches(0.9), Inches(6.15), W - Inches(1.8), Inches(0.4),
              'Based on: Chen, Liu, Wang, Yin — "Hyperparameter Tuning is All You Need for LISTA" · NeurIPS 2021',
              size=13, color=RGBColor(0x88, 0x99, 0xAA), italic=True)


def slide_02_roadmap(prs):
    slide = _add_slide(prs)
    _slide_bg(slide, BG)
    _top_bar(slide)
    _bottom_bar(slide)
    _add_title(slide, 'Seminar Roadmap')
    _page_number(slide, 2)

    steps = [
        ('1', 'Background', 'Motivation, LASSO,\nISTA / FISTA', False),
        ('2', 'Unrolling', 'LISTA, ALISTA,\nHyperLISTA', False),
        ('3', 'Our Study', 'Scalar variants +\ntraining methods', True),
        ('4', 'Results', 'Sparse recovery &\ndesign space', False),
        ('5', 'Images', 'MNIST & FashionMNIST\nCS stress tests', False),
    ]
    sx = Inches(0.55)
    sw = (W - Inches(1.1)) / 5 - Inches(0.12)
    sy = Inches(1.3)
    sh = Inches(2.0)

    for i, (num, title, sub, active) in enumerate(steps):
        cx = sx + i * (sw + Inches(0.12))
        bg_c = RGBColor(0xEA, 0xF5, 0xFF) if active else WHITE
        bdr_c = ACCENT if active else BORDER
        shape = slide.shapes.add_shape(1, cx, sy, sw, sh)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_c
        shape.line.color.rgb = bdr_c
        shape.line.width = Pt(1.5 if active else 0.75)

        # circle with number
        circ_r = Inches(0.22)
        circ_x = cx + sw / 2 - circ_r
        circ_y = sy + Inches(0.2)
        circ = slide.shapes.add_shape(9, circ_x, circ_y, circ_r * 2, circ_r * 2)  # oval
        circ.fill.solid()
        circ.fill.fore_color.rgb = ACCENT if active else BGU_BLUE
        circ.line.fill.background()
        txb_n = slide.shapes.add_textbox(circ_x, circ_y, circ_r * 2, circ_r * 2)
        p = txb_n.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = num
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = WHITE
        run.font.name = 'Calibri'

        tc = ACCENT if active else BGU_BLUE
        _text_box(slide, cx + Inches(0.08), sy + Inches(0.75), sw - Inches(0.16), Inches(0.4),
                  title, size=15, bold=True, color=tc, align=PP_ALIGN.CENTER)
        _text_box(slide, cx + Inches(0.08), sy + Inches(1.2), sw - Inches(0.16), Inches(0.65),
                  sub, size=12, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

    # Goal text
    _text_box(slide, Inches(0.55), Inches(3.55), W - Inches(1.1), Inches(0.4),
              'Goal:  systematically compare how much learning is needed when moving from ideal sparse signals to real compressed-sensing images.',
              size=16, color=TEXT)

    # 3 info cards
    cards = [
        ('Synthetic signals', 'm=250, n=500, s=50\nGaussian A, noiseless'),
        ('Trained models', 'LISTA, ALISTA, HyperLISTA\n+ 3 scalar designs (ours)'),
        ('Image CS', 'MNIST & FashionMNIST pixel domain\nratios ∈ {0.25, 0.5}'),
    ]
    cw = (W - Inches(1.1)) / 3 - Inches(0.12)
    cy = Inches(4.15)
    ch = Inches(1.1)
    for i, (ct, cb) in enumerate(cards):
        cx2 = Inches(0.55) + i * (cw + Inches(0.12))
        _card(slide, cx2, cy, cw, ch, ct, cb.split('\n'),
              left_accent=ACCENT, bg=RGBColor(0xE8, 0xF5, 0xFF),
              border_color=RGBColor(0xC0, 0xE0, 0xFF))


def slide_03_motivation(prs):
    slide = _add_slide(prs)
    _slide_bg(slide, BG)
    _top_bar(slide)
    _bottom_bar(slide)
    _add_title(slide, 'Motivation & Problem Formulation')
    _page_number(slide, 3)

    # Left column: applications table
    lx = Inches(0.55)
    ly = Inches(1.25)
    lw = Inches(6.5)

    _text_box(slide, lx, ly, lw, Inches(0.35), 'Sparse recovery appears everywhere:', size=17, bold=True, color=BGU_BLUE)

    rows = [
        ['Application', 'Measurements', 'Goal'],
        ['MRI acceleration', 'k-space  (m ≪ n)', 'Image pixels'],
        ['Radar', 'Compressed snapshots', 'Reflectivity'],
        ['Hyperspectral', 'Spectral bands', 'Spatial cube'],
        ['Compressed sensing', 'Random projections', 'Original signal'],
    ]
    cws = [Inches(2.2), Inches(2.4), Inches(1.7)]
    _add_table(slide, lx, ly + Inches(0.4), lw, rows, cws, row_size=14)

    _text_box(slide, lx, ly + Inches(2.65), lw, Inches(0.45),
              'Why tractable?  Natural signals have only s ≪ n non-zero components.', size=15, color=TEXT)

    _text_box(slide, lx, ly + Inches(3.15), lw, Inches(0.35), 'LASSO — the optimization backbone:', size=17, bold=True, color=BGU_BLUE)

    eq1 = render_eq(r'\hat{x} = \arg\min_{x}\; \tfrac{1}{2}\|b - Ax\|_2^2 + \lambda\|x\|_1', fontsize=22)
    _add_eq_image(slide, eq1, lx + Inches(0.2), ly + Inches(3.55), Inches(4.8))

    # Right column: dark card
    rx = Inches(7.3)
    ry = Inches(1.25)
    rw = Inches(5.5)
    rh = Inches(3.2)
    bg_shape = slide.shapes.add_shape(1, rx, ry, rw, rh)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = BGU_BLUE
    bg_shape.line.fill.background()

    _text_box(slide, rx + Inches(0.2), ry + Inches(0.2), rw - Inches(0.4), Inches(0.4),
              'Measurement model', size=16, bold=True, color=GOLD)

    eq2 = render_eq(r'b = Ax^* + \varepsilon', fontsize=26, color='white')
    _add_eq_image(slide, eq2, rx + Inches(0.5), ry + Inches(0.7), Inches(3.0))

    lines = [
        'A ∈ ℝ^{m×n},   m ≪ n',
        'x*: s-sparse signal',
        'ε: additive Gaussian noise',
    ]
    for i, ln in enumerate(lines):
        _text_box(slide, rx + Inches(0.2), ry + Inches(1.6) + i * Inches(0.42),
                  rw - Inches(0.4), Inches(0.4), ln, size=15, color=WHITE)

    # Metrics note
    _card(slide, rx, ry + Inches(3.4), rw, Inches(1.0),
          None,
          ['Metrics:  NMSE (dB) for synthetic signals;',
           'NMSE + PSNR + SSIM for image CS.'],
          left_accent=ACCENT, bg=RGBColor(0xE8, 0xF5, 0xFF),
          border_color=RGBColor(0xC0, 0xE0, 0xFF), body_size=15)


def slide_04_baselines(prs):
    slide = _add_slide(prs)
    _slide_bg(slide, BG)
    _top_bar(slide)
    _bottom_bar(slide)
    _add_title(slide, 'Classical Baselines: ISTA and FISTA')
    _page_number(slide, 4)

    y = Inches(1.3)
    _text_box(slide, Inches(0.55), y, W - Inches(1.1), Inches(0.35),
              'ISTA  —  proximal gradient descent on the LASSO objective:', size=17, bold=True, color=BGU_BLUE)
    eq = render_eq(r'x^{(k+1)} = \mathcal{S}_{\theta}\!\left(x^{(k)} + \tfrac{1}{L} A^T(b - Ax^{(k)})\right), \quad L = \|A\|_2^2,\quad \theta = \tfrac{\lambda}{L}', fontsize=20)
    _add_eq_image(slide, eq, Inches(0.9), y + Inches(0.4), Inches(10.5))

    y2 = Inches(2.55)
    _text_box(slide, Inches(0.55), y2, W - Inches(1.1), Inches(0.35),
              'FISTA  —  ISTA with Nesterov momentum,  O(1/k²) convergence vs O(1/k):', size=17, bold=True, color=BGU_BLUE)
    eq2 = render_eq(r'y^{(k)} = x^{(k)} + \tfrac{t_k-1}{t_{k+1}}\bigl(x^{(k)}-x^{(k-1)}\bigr),\quad x^{(k+1)} = \mathcal{S}_\theta\!\left(y^{(k)} + \tfrac{1}{L}A^T(b - Ay^{(k)})\right)', fontsize=20)
    _add_eq_image(slide, eq2, Inches(0.9), y2 + Inches(0.4), Inches(11.0))

    # Results table + callout
    rows = [
        ['Method', 'NMSE @ K=16', 'Params'],
        ['ISTA', '−5.3 dB', '0'],
        ['FISTA', '−11.0 dB', '0'],
    ]
    cws = [Inches(3.0), Inches(2.5), Inches(1.8)]
    _add_table(slide, Inches(0.55), Inches(4.15), Inches(7.5), rows, cws, row_size=15)

    _card(slide, Inches(8.3), Inches(4.15), Inches(4.7), Inches(1.3),
          None,
          ['At K=16 iterations both methods are far from converged.',
           'ISTA needs 200+ iterations for −20 dB.',
           '',
           'Can we make 16 steps count?'],
          left_accent=ACCENT, bg=RGBColor(0xE8, 0xF5, 0xFF),
          border_color=RGBColor(0xC0, 0xE0, 0xFF), body_size=15)


def slide_05_unrolling(prs):
    slide = _add_slide(prs)
    _slide_bg(slide, BG)
    _top_bar(slide)
    _bottom_bar(slide)
    _add_title(slide, 'Deep Unrolling: From Iterations to Layers')
    _page_number(slide, 5)

    _text_box(slide, Inches(0.55), Inches(1.3), W - Inches(1.1), Inches(0.38),
              'Key idea (Gregor & LeCun, 2010):  unfold K ISTA iterations into a K-layer network and learn the update parameters from data.',
              size=17, color=TEXT)

    # Two cards
    card_y = Inches(1.85)
    card_h = Inches(3.2)
    card_w = (W - Inches(1.35)) / 2

    # Left card — Classical ISTA
    lx = Inches(0.55)
    shape_l = slide.shapes.add_shape(1, lx, card_y, card_w, card_h)
    shape_l.fill.solid(); shape_l.fill.fore_color.rgb = WHITE
    shape_l.line.color.rgb = BORDER; shape_l.line.width = Pt(0.75)

    _text_box(slide, lx + Inches(0.2), card_y + Inches(0.15), card_w - Inches(0.4), Inches(0.38),
              'Classical ISTA  —  parameters fixed', size=16, bold=True, color=BGU_BLUE)

    eq_l1 = render_eq(r'z^{(k)} = x^{(k)} + \tfrac{1}{L} A^T(b - Ax^{(k)})', fontsize=19)
    _add_eq_image(slide, eq_l1, lx + Inches(0.15), card_y + Inches(0.65), Inches(5.6))

    eq_l2 = render_eq(r'x^{(k+1)} = \mathcal{S}_{\theta}(z^{(k)})', fontsize=19)
    _add_eq_image(slide, eq_l2, lx + Inches(0.15), card_y + Inches(1.35), Inches(4.0))

    _text_box(slide, lx + Inches(0.2), card_y + Inches(2.1), card_w - Inches(0.4), Inches(0.4),
              'Parameters 1/L and θ = λ/L  are fixed constants.', size=14, color=TEXT_LIGHT)
    _rect(slide, lx + Inches(0.2), card_y + Inches(2.6), card_w - Inches(0.4), Inches(0.38),
          RGBColor(0xEE, 0xF4, 0xFB))
    _text_box(slide, lx + Inches(0.25), card_y + Inches(2.62), card_w - Inches(0.5), Inches(0.35),
              'Fixed  —  no learning', size=14, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

    # Right card — LISTA
    rx = lx + card_w + Inches(0.25)
    shape_r = slide.shapes.add_shape(1, rx, card_y, card_w, card_h)
    shape_r.fill.solid(); shape_r.fill.fore_color.rgb = RGBColor(0xF0, 0xF8, 0xFF)
    shape_r.line.color.rgb = ACCENT; shape_r.line.width = Pt(1.5)

    _text_box(slide, rx + Inches(0.2), card_y + Inches(0.15), card_w - Inches(0.4), Inches(0.38),
              'Unrolled LISTA  —  parameters learned', size=16, bold=True, color=BGU_BLUE)

    eq_r1 = render_eq(r'z^{(k)} = W_x^{(k)} x^{(k)} + W_y^{(k)} b', fontsize=19)
    _add_eq_image(slide, eq_r1, rx + Inches(0.15), card_y + Inches(0.65), Inches(5.0))

    eq_r2 = render_eq(r'x^{(k+1)} = \mathcal{S}_{\theta_k}(z^{(k)})', fontsize=19)
    _add_eq_image(slide, eq_r2, rx + Inches(0.15), card_y + Inches(1.35), Inches(4.0))

    _text_box(slide, rx + Inches(0.2), card_y + Inches(2.1), card_w - Inches(0.4), Inches(0.4),
              'W_x(k), W_y(k), θ_k  are learned per layer via BPTT.', size=14, color=TEXT_LIGHT)
    _rect(slide, rx + Inches(0.2), card_y + Inches(2.6), card_w - Inches(0.4), Inches(0.38),
          RGBColor(0x00, 0xA3, 0xE0))
    _text_box(slide, rx + Inches(0.25), card_y + Inches(2.62), card_w - Inches(0.5), Inches(0.35),
              '~6 M parameters  —  learned end-to-end', size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # MBDL principle callout
    _card(slide, Inches(0.55), Inches(5.25), W - Inches(1.1), Inches(1.1),
          None,
          ['MBDL principle:  preserve the known algorithmic structure (gradient step + soft-threshold);',
           'learn only what improves finite-depth reconstruction. Training minimizes  ‖x(K) − x*‖²  via BPTT.'],
          left_accent=ACCENT, bg=RGBColor(0xE8, 0xF5, 0xFF),
          border_color=RGBColor(0xC0, 0xE0, 0xFF), body_size=15)


def slide_06_mbdl_ladder(prs):
    slide = _add_slide(prs)
    _slide_bg(slide, BG)
    _top_bar(slide)
    _bottom_bar(slide)
    _add_title(slide, 'The MBDL Ladder: LISTA → ALISTA → HyperLISTA')
    _page_number(slide, 6)

    lx = Inches(0.55)
    rw = Inches(4.5)
    rx = W - rw - Inches(0.55)
    lw = W - rw - Inches(1.5)

    # LISTA
    _text_box(slide, lx, Inches(1.3), lw, Inches(0.35),
              'LISTA  —  learn all matrices per layer,  ~6 M params:', size=16, bold=True, color=BGU_BLUE)
    eq1 = render_eq(r'x^{(k+1)} = \mathcal{S}_{\theta_k}\!\left(W_x^{(k)} x^{(k)} + W_y^{(k)} b\right)', fontsize=18)
    _add_eq_image(slide, eq1, lx, Inches(1.7), Inches(7.0))

    # ALISTA
    _text_box(slide, lx, Inches(2.55), lw, Inches(0.35),
              'ALISTA  —  W analytic, learn step sizes and thresholds,  32 params:', size=16, bold=True, color=BGU_BLUE)
    eq2a = render_eq(r'W^* = (AA^T)^{-1}A \;\text{(normalized)}', fontsize=17)
    _add_eq_image(slide, eq2a, lx, Inches(2.95), Inches(5.5))
    eq2b = render_eq(r'x^{(k+1)} = \mathcal{S}_{\theta_k}\!\left(x^{(k)} + \gamma_k W^{*T}(b - Ax^{(k)})\right)', fontsize=17)
    _add_eq_image(slide, eq2b, lx, Inches(3.4), Inches(7.5))

    # HyperLISTA
    _text_box(slide, lx, Inches(4.2), lw, Inches(0.35),
              'HyperLISTA  —  adaptive θ/β/p from residuals, 3 params (grid search, zero backprop):', size=16, bold=True, color=BGU_BLUE)
    eq3a = render_eq(r'\theta^{(k)} = c_1\mu\|A^+(Ax^{(k)}\!-b)\|_1,\quad \beta^{(k)} = c_2\mu\|x^{(k)}\|_0', fontsize=15)
    _add_eq_image(slide, eq3a, lx, Inches(4.6), Inches(7.0))
    eq3b = render_eq(r'p^{(k)} = c_3\log\tfrac{\|A^+b\|_1}{\|A^+(Ax^{(k)}-b)\|_1}', fontsize=15)
    _add_eq_image(slide, eq3b, lx, Inches(5.1), Inches(5.5))

    # Right: table
    rows = [
        ['Model', 'Params', 'NMSE'],
        ['ISTA', '0', '−5.3 dB'],
        ['FISTA', '0', '−11.0 dB'],
        ['LISTA', '6 M', '−20.3 dB'],
        ['ALISTA', '32', '−29.8 dB'],
        ['HyperLISTA', '3', '−61.4 dB'],
    ]
    cws2 = [Inches(2.0), Inches(1.1), Inches(1.2)]
    _add_table(slide, rx, Inches(1.3), rw, rows, cws2,
               highlight_rows={5}, row_size=15)

    # result box
    rb = slide.shapes.add_shape(1, rx, Inches(4.15), rw, Inches(1.35))
    rb.fill.solid(); rb.fill.fore_color.rgb = BGU_BLUE
    rb.line.fill.background()
    txb = slide.shapes.add_textbox(rx + Inches(0.15), Inches(4.2), rw - Inches(0.3), Inches(1.25))
    tf = txb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = 'More structure → fewer params → better performance'
    run.font.size = Pt(16); run.font.bold = True; run.font.color.rgb = WHITE; run.font.name = 'Calibri'
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = 'This is the MBDL principle in action.'
    run2.font.size = Pt(14); run2.font.color.rgb = GOLD; run2.font.name = 'Calibri'


def slide_07_design(prs):
    slide = _add_slide(prs)
    _slide_bg(slide, BG)
    _top_bar(slide)
    _bottom_bar(slide)
    _add_title(slide, 'Our Design Space Study: What Should We Learn?')
    _page_number(slide, 7)

    _text_box(slide, Inches(0.55), Inches(1.3), W - Inches(1.1), Inches(0.38),
              'Research question:  LISTA learns 6 M parameters; HyperLISTA uses 3 scalars. What is the minimal useful set of learnable components?',
              size=17, color=TEXT)

    _text_box(slide, Inches(0.55), Inches(1.78), W - Inches(1.1), Inches(0.35),
              'The ISTA gradient-correction update has two free scalar sequences:', size=16, color=TEXT)

    eq = render_eq(r'x^{(k+1)} = \mathcal{S}_{\theta_k}\!\left(x^{(k)} + \gamma_k A^T(b - Ax^{(k)})\right)', fontsize=22)
    _add_eq_image(slide, eq, Inches(1.5), Inches(2.2), Inches(9.0))

    _text_box(slide, Inches(0.55), Inches(3.05), W - Inches(1.1), Inches(0.35),
              'We designed three models, each learning a different subset:', size=16, color=TEXT)

    # 3 cards
    cards = [
        ('ThresholdISTA', BGU_BLUE, BORDER, WHITE, [
            '16 params',
            'Learns  θ₁, …, θ_K  only',
            'Step fixed at  γ = 1/L',
            'Init:  θ₀ = λ/L',
        ]),
        ('StepISTA', ACCENT, ACCENT, RGBColor(0xE8, 0xF5, 0xFF), [
            '16 params',
            'Learns  γ₁, …, γ_K  only',
            'Threshold coupled:  θ_k = λγ_k',
            'Init:  γ₀ = 1/L',
        ]),
        ('StepThresholdISTA', GOLD, GOLD, RGBColor(0xFF, 0xF8, 0xE8), [
            '32 params',
            'Learns both  γ_k  and  θ_k',
            'independently per layer',
            'Init:  ISTA values',
        ]),
    ]
    cw = (W - Inches(1.35)) / 3
    for i, (title, hdr, bdr, bg_c, lines) in enumerate(cards):
        cx = Inches(0.55) + i * (cw + Inches(0.075))
        shape = slide.shapes.add_shape(1, cx, Inches(3.5), cw, Inches(2.7))
        shape.fill.solid(); shape.fill.fore_color.rgb = bg_c
        shape.line.color.rgb = bdr; shape.line.width = Pt(1.5)

        _text_box(slide, cx + Inches(0.15), Inches(3.6), cw - Inches(0.3), Inches(0.4),
                  title, size=17, bold=True, color=hdr)

        badge = slide.shapes.add_shape(1, cx + Inches(0.15), Inches(4.05), Inches(1.1), Inches(0.3))
        badge.fill.solid(); badge.fill.fore_color.rgb = hdr
        badge.line.fill.background()
        _text_box(slide, cx + Inches(0.16), Inches(4.06), Inches(1.08), Inches(0.28),
                  lines[0], size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        for j, ln in enumerate(lines[1:]):
            _text_box(slide, cx + Inches(0.15), Inches(4.45) + j * Inches(0.38),
                      cw - Inches(0.3), Inches(0.36), ln, size=14, color=TEXT)

    _text_box(slide, Inches(0.55), Inches(6.35), W - Inches(1.1), Inches(0.3),
              'All parameters softplus-reparameterized (positivity). All initialized to reproduce ISTA exactly before training.',
              size=13, color=TEXT_LIGHT, italic=True)


def slide_08_scalar_results(prs):
    slide = _add_slide(prs)
    _slide_bg(slide, BG)
    _top_bar(slide)
    _bottom_bar(slide)
    _add_title(slide, 'Our Results: Scalar Models')
    _page_number(slide, 8)

    fig_path = os.path.join(ASSETS, 'scalar_models.png')
    _add_image(slide, fig_path, Inches(0.35), Inches(1.2), Inches(7.4))

    rows = [
        ['Model', 'NMSE', 'Params'],
        ['ISTA', '−5.3 dB', '0'],
        ['FISTA', '−11.0 dB', '0'],
        ['ThresholdISTA', '−7.1 dB', '16'],
        ['StepISTA', '−20.0 dB', '16'],
        ['StepThresholdISTA', '−23.6 dB', '32'],
        ['LISTA-L1  (ref.)', '−20.3 dB', '6 M'],
    ]
    cws = [Inches(2.35), Inches(1.3), Inches(1.0)]
    _add_table(slide, Inches(7.95), Inches(1.2), Inches(5.0), rows, cws,
               highlight_rows={5}, highlight2_rows={4}, row_size=14)

    _card(slide, Inches(7.95), Inches(4.25), Inches(5.0), Inches(2.15),
          'Key finding', [
              'Learning γ_k alone (16 scalars) already',
              'matches full LISTA (6 M parameters).',
              '',
              'Decoupling θ_k adds 3 dB at zero extra cost.',
              '',
              'StepThresholdISTA beats LISTA-L1',
              'with 187,000× fewer parameters.',
          ],
          title_color=RGBColor(0x7A, 0x4F, 0x00),
          bg=RGBColor(0xFF, 0xF8, 0xE8),
          border_color=GOLD, left_accent=GOLD, body_size=13)


def slide_09_training(prs):
    slide = _add_slide(prs)
    _slide_bg(slide, BG)
    _top_bar(slide)
    _bottom_bar(slide)
    _add_title(slide, 'Training Methods for Unfolded Optimizers')
    _page_number(slide, 9)

    _text_box(slide, Inches(0.55), Inches(1.28), W - Inches(1.1), Inches(0.35),
              'Architecture fixed (LISTA). Only the training strategy changes.', size=17, bold=True, color=BGU_BLUE)

    # 3 method cards
    methods = [
        ('L1 — End-to-end', r'\mathcal{L} = \|x^{(K)} - x^*\|^2', 'Supervise final output only.\nStandard BPTT.'),
        ('L2 — Deep Supervision', r'\mathcal{L} = \|x^{(K)}-x^*\|^2 + \lambda\sum_k w_k\|x^{(k)}-x^*\|^2', 'Supervise every\nintermediate layer.'),
        ('L3 — Greedy Layerwise', r'\min_{\theta_k} \|x^{(k)} - x^*\|^2', 'Train layer k while\nlayers 1…k−1 frozen.'),
    ]
    cw = (W - Inches(1.35)) / 3
    eq_widths = [Inches(3.1), cw - Inches(0.2), Inches(2.5)]
    for i, (title, latex, desc) in enumerate(methods):
        cx = Inches(0.55) + i * (cw + Inches(0.075))
        shape = slide.shapes.add_shape(1, cx, Inches(1.75), cw, Inches(2.3))
        shape.fill.solid(); shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = BORDER; shape.line.width = Pt(0.75)
        _text_box(slide, cx + Inches(0.15), Inches(1.83), cw - Inches(0.3), Inches(0.38),
                  title, size=16, bold=True, color=BGU_BLUE)
        eq = render_eq(latex, fontsize=16)
        _add_eq_image(slide, eq, cx + Inches(0.1), Inches(2.3), eq_widths[i])
        _text_box(slide, cx + Inches(0.15), Inches(3.25), cw - Inches(0.3), Inches(0.7),
                  desc, size=14, color=TEXT)

    # Tied vs untied table
    _text_box(slide, Inches(0.55), Inches(4.25), W - Inches(1.1), Inches(0.35),
              'Tied vs. Untied weights:', size=17, bold=True, color=BGU_BLUE)

    rows2 = [
        ['', 'Untied LISTA', 'Tied LISTA'],
        ['Weights', 'Independent per layer', 'Single shared (Wx, Wy, θ)'],
        ['Params', '~6 M', '375 K'],
        ['Style', 'Fully flexible', 'RNN-like'],
        ['Effect', 'Hard to optimize', 'Smoother loss landscape'],
    ]
    cws2 = [Inches(1.3), Inches(3.6), Inches(3.6)]
    _add_table(slide, Inches(0.55), Inches(4.65), Inches(8.5), rows2, cws2, row_size=14)

    _card(slide, Inches(9.4), Inches(4.65), Inches(3.6), Inches(1.75),
          None,
          ['Weight tying = implicit regularization.',
           '',
           'LISTA-Tied-L1:  −19.1 dB @ 375 K',
           'LISTA-L1:  −20.3 dB @ 6 M'],
          left_accent=ACCENT, bg=RGBColor(0xE8, 0xF5, 0xFF),
          border_color=RGBColor(0xC0, 0xE0, 0xFF), body_size=12)


def slide_10_sparse_results(prs):
    slide = _add_slide(prs)
    _slide_bg(slide, BG)
    _top_bar(slide)
    _bottom_bar(slide)
    _add_title(slide, 'Synthetic Sparse Recovery Results')
    _page_number(slide, 10)

    fig_path = os.path.join(ASSETS, 'main_models.png')
    _add_image(slide, fig_path, Inches(0.35), Inches(1.2), Inches(7.4))

    rows = [
        ['Method', 'NMSE', 'Params'],
        ['ISTA', '−5.3 dB', '0'],
        ['FISTA', '−11.0 dB', '0'],
        ['LISTA-L1', '−20.3 dB', '6 M'],
        ['LISTA-L2', '−21.9 dB', '6 M'],
        ['LISTA-L3', '−16.6 dB', '6 M'],
        ['LISTA-Tied-L1', '−19.1 dB', '375 K'],
        ['ALISTA', '−29.8 dB', '32'],
        ['HyperLISTA', '−61.4 dB', '3'],
    ]
    cws = [Inches(2.3), Inches(1.4), Inches(1.0)]
    _add_table(slide, Inches(7.95), Inches(1.2), Inches(5.0), rows, cws,
               highlight_rows={8}, highlight2_rows={7}, row_size=13)

    _card(slide, Inches(7.95), Inches(4.65), Inches(5.0), Inches(1.8),
          None,
          ['L3 plateau:  NMSE flattens ~layer 8.',
           'Greedy training gives no incentive to improve',
           'beyond what prior layers solved.',
           '',
           'LISTA-Tied > LISTA-Untied despite 16× fewer params.'],
          left_accent=ACCENT, bg=RGBColor(0xE8, 0xF5, 0xFF),
          border_color=RGBColor(0xC0, 0xE0, 0xFF), body_size=13)


def slide_11_design_space(prs):
    slide = _add_slide(prs)
    _slide_bg(slide, BG)
    _top_bar(slide)
    _bottom_bar(slide)
    _add_title(slide, 'Design-Space Map')
    _page_number(slide, 11)

    fig_path = os.path.join(ASSETS, 'param_tradeoff_full.png')
    _add_image(slide, fig_path, Inches(0.35), Inches(1.2), Inches(7.0))

    rows = [
        ['Model', 'Params', 'NMSE', 'What drives it'],
        ['ThresholdISTA', '16', '−7.1 dB', 'learned θ_k'],
        ['StepISTA', '16', '−20.0 dB', 'learned γ_k'],
        ['StepThresholdISTA', '32', '−23.6 dB', 'γ_k + θ_k'],
        ['LISTA-L1', '6 M', '−20.3 dB', 'full matrices'],
        ['ALISTA', '32', '−29.8 dB', 'analytic W'],
        ['HyperLISTA', '3', '−61.4 dB', 'adaptive residuals'],
    ]
    cws2 = [Inches(2.35), Inches(0.8), Inches(1.1), Inches(1.7)]
    _add_table(slide, Inches(7.55), Inches(1.2), Inches(5.4), rows, cws2,
               highlight_rows={3, 6}, row_size=13)

    # result box
    rb = slide.shapes.add_shape(1, Inches(7.55), Inches(4.75), Inches(5.4), Inches(1.45))
    rb.fill.solid(); rb.fill.fore_color.rgb = BGU_BLUE; rb.line.fill.background()
    txb = slide.shapes.add_textbox(Inches(7.75), Inches(4.82), Inches(5.0), Inches(1.3))
    tf = txb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; run = p.add_run()
    run.text = 'StepThresholdISTA  (32 scalars)'; run.font.size = Pt(16)
    run.font.bold = True; run.font.color.rgb = GOLD; run.font.name = 'Calibri'
    p2 = tf.add_paragraph(); run2 = p2.add_run()
    run2.text = 'beats LISTA-L1  (6 million matrices).'; run2.font.size = Pt(16)
    run2.font.bold = True; run2.font.color.rgb = WHITE; run2.font.name = 'Calibri'
    p3 = tf.add_paragraph(); run3 = p3.add_run()
    run3.text = 'The step-size schedule γ_k is the bottleneck.'; run3.font.size = Pt(14)
    run3.font.color.rgb = RGBColor(0xAA, 0xCC, 0xFF); run3.font.name = 'Calibri'


def slide_12_images(prs):
    slide = _add_slide(prs)
    _slide_bg(slide, BG)
    _top_bar(slide)
    _bottom_bar(slide)
    _add_title(slide, 'Real Image CS: FashionMNIST Pixel Domain')
    _page_number(slide, 12)

    _text_box(slide, Inches(0.55), Inches(1.3), W - Inches(1.1), Inches(0.38),
              'Setup:  flatten 28×28 image to x ∈ [0,1]^{784},  measure  b = Ax,  ratio m/784 ∈ {0.25, 0.5}.  FashionMNIST: ~51.5% of pixels exactly zero.',
              size=16, color=TEXT)

    rows = [
        ['Method', 'NMSE', 'PSNR', 'SSIM', 'Params'],
        ['ISTA', '−1.2 dB', '8.9', '0.13', '0'],
        ['FISTA', '−1.2 dB', '8.9', '0.14', '0'],
        ['StepThresholdISTA', '−1.2 dB', '8.9', '0.13', '32'],
        ['ALISTA', '−0.9 dB', '8.4', '0.12', '32'],
        ['HyperLISTA', '−0.3 dB', '7.7', '0.09', '3'],
        ['LISTA-L2', '−13.7 dB', '22.3', '0.82', '12 M'],
    ]
    cws = [Inches(2.5), Inches(1.3), Inches(0.9), Inches(0.9), Inches(1.1)]
    _add_table(slide, Inches(0.55), Inches(1.85), Inches(7.0), rows, cws,
               highlight_rows={6}, row_size=14)

    _text_box(slide, Inches(0.55), Inches(5.0), Inches(7.0), Inches(0.3),
              'ratio = 0.25', size=13, color=TEXT_LIGHT, italic=True)

    # Two explanation cards
    _card(slide, Inches(7.75), Inches(1.85), Inches(5.2), Inches(1.75),
          'Signal model mismatch', [
              'All structured methods assume i.i.d.',
              'Gaussian sparse signals.',
              'FashionMNIST pixels are bounded [0,1],',
              'spatially correlated, and only soft-sparse.',
          ],
          title_color=RGBColor(0x7A, 0x4F, 0x00),
          bg=RGBColor(0xFF, 0xF8, 0xE8), border_color=GOLD,
          left_accent=GOLD, body_size=13)

    _card(slide, Inches(7.75), Inches(3.75), Inches(5.2), Inches(1.75),
          'Why LISTA wins', [
              'LISTA-L2 is assumption-free —',
              'it learns the actual image distribution',
              'from 51 K training images.',
              '',
              'Wrong prior is worse than no prior.',
          ],
          left_accent=ACCENT, bg=RGBColor(0xE8, 0xF5, 0xFF),
          border_color=RGBColor(0xC0, 0xE0, 0xFF), body_size=13)


def slide_13_robustness(prs):
    slide = _add_slide(prs)
    _slide_bg(slide, BG)
    _top_bar(slide)
    _bottom_bar(slide)
    _add_title(slide, 'Robustness Across Datasets and Measurement Ratios')
    _page_number(slide, 13)

    _text_box(slide, Inches(0.55), Inches(1.25), W - Inches(1.1), Inches(0.38),
              'SSIM across {MNIST, FashionMNIST} × ratio {0.25, 0.5}:  harder data and fewer measurements widen the gap.',
              size=16, color=TEXT)

    _add_image(slide, os.path.join(ASSETS, 'ssim_heatmap.png'),
               Inches(0.45), Inches(1.8), Inches(6.4))
    _add_image(slide, os.path.join(ASSETS, 'ssim_bars.png'),
               Inches(7.05), Inches(1.85), Inches(5.9))

    cards2 = [
        ('Trend', ACCENT,
         'MNIST @ 0.5 is easy for everyone; FashionMNIST @ 0.25 breaks every structured model. LISTA-L2 stays ≥ 0.82 SSIM everywhere.'),
        ('Interpretation', GOLD,
         'The sparse-Gaussian prior degrades gracefully on MNIST (~78% zeros) but fails on FashionMNIST (~51% zeros, textured).'),
    ]
    cw2 = (W - Inches(1.25)) / 2
    for i, (ct, ac, cb) in enumerate(cards2):
        cx = Inches(0.55) + i * (cw2 + Inches(0.15))
        _card(slide, cx, Inches(5.55), cw2, Inches(1.5),
              ct, [cb],
              title_color=BGU_BLUE, left_accent=ac,
              bg=RGBColor(0xE8, 0xF5, 0xFF) if ac == ACCENT else RGBColor(0xFF, 0xF8, 0xE8),
              border_color=RGBColor(0xC0, 0xE0, 0xFF) if ac == ACCENT else GOLD,
              body_size=13)


def slide_14_stress(prs):
    slide = _add_slide(prs)
    _slide_bg(slide, BG)
    _top_bar(slide)
    _bottom_bar(slide)
    _add_title(slide, 'Visual Stress Tests  (ratio = 0.25)')
    _page_number(slide, 14)

    # Two reconstruction grids side by side (rows: GT, ISTA, FISTA, StepThr, LISTA-L2, ALISTA, HyperLISTA)
    img_w = Inches(5.15)
    _add_image(slide, os.path.join(ASSETS, 'mnist_grid_025.png'),
               Inches(0.85), Inches(1.2), img_w)
    _add_image(slide, os.path.join(ASSETS, 'fmnist_grid_025.png'),
               Inches(7.35), Inches(1.2), img_w)

    _card(slide, Inches(0.55), Inches(6.1), W - Inches(1.1), Inches(1.05),
          None,
          ['LISTA-L2 recovers clean structure on both datasets. ALISTA / HyperLISTA zero out valid signal on FashionMNIST —',
           'the Gaussian-sparse threshold is miscalibrated for textured images. Wrong prior is worse than no prior.'],
          left_accent=ACCENT, bg=RGBColor(0xE8, 0xF5, 0xFF),
          border_color=RGBColor(0xC0, 0xE0, 0xFF), body_size=14)


def slide_14_conclusions(prs):
    slide = _add_slide(prs)
    _slide_bg(slide, BG)
    _top_bar(slide)
    _bottom_bar(slide)
    _add_title(slide, 'Conclusions')
    _page_number(slide, 15)

    findings = [
        ('1', 'Structure wins when assumptions match.',
         'HyperLISTA\'s 3 scalars beat 6 M learned parameters by 41 dB on synthetic sparse data.'),
        ('2', 'Step sizes matter most.',
         'StepISTA (16 scalars) matches LISTA. Decoupling θ_k adds 3 dB. Matrix learning is largely redundant.'),
        ('3', 'Weight tying is implicit regularization.',
         'LISTA-Tied matches LISTA-Untied (−19.1 vs −20.3 dB) with 16× fewer parameters.'),
        ('4', 'Wrong prior actively hurts.',
         'ALISTA / HyperLISTA score below ISTA on real images. LISTA wins by making no structural assumptions.'),
    ]
    cw = (W - Inches(1.1)) / 2 - Inches(0.08)
    for i, (num, title, body) in enumerate(findings):
        col = i % 2
        row = i // 2
        cx = Inches(0.55) + col * (cw + Inches(0.16))
        cy = Inches(1.3) + row * Inches(1.55)

        shape = slide.shapes.add_shape(1, cx, cy, cw, Inches(1.4))
        shape.fill.solid(); shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = BORDER; shape.line.width = Pt(0.75)

        # Number badge
        circ = slide.shapes.add_shape(9, cx + Inches(0.15), cy + Inches(0.15), Inches(0.42), Inches(0.42))
        circ.fill.solid(); circ.fill.fore_color.rgb = BGU_BLUE; circ.line.fill.background()
        _text_box(slide, cx + Inches(0.15), cy + Inches(0.15), Inches(0.42), Inches(0.42),
                  num, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        _text_box(slide, cx + Inches(0.7), cy + Inches(0.15), cw - Inches(0.85), Inches(0.4),
                  title, size=16, bold=True, color=BGU_BLUE)
        _text_box(slide, cx + Inches(0.7), cy + Inches(0.6), cw - Inches(0.85), Inches(0.72),
                  body, size=13, color=TEXT)

    # Result box
    rb = slide.shapes.add_shape(1, Inches(0.55), Inches(4.55), Inches(6.0), Inches(1.65))
    rb.fill.solid(); rb.fill.fore_color.rgb = BGU_BLUE; rb.line.fill.background()
    txb = slide.shapes.add_textbox(Inches(0.75), Inches(4.62), Inches(5.6), Inches(1.5))
    tf = txb.text_frame; tf.word_wrap = True
    for line, col, sz, bold in [
        ('MBDL hierarchy on synthetic data:', GOLD, 15, True),
        ('ISTA → LISTA → ALISTA → HyperLISTA', WHITE, 14, False),
        ('0  →  6M  →  32  →  3  params', ACCENT, 14, False),
        ('−5  →  −20  →  −30  →  −61  dB', ACCENT, 14, True),
    ]:
        p = tf.paragraphs[0] if line == 'MBDL hierarchy on synthetic data:' else tf.add_paragraph()
        run = p.add_run(); run.text = line
        run.font.size = Pt(sz); run.font.color.rgb = col
        run.font.bold = bold; run.font.name = 'Calibri'

    # Future work
    _card(slide, Inches(6.75), Inches(4.55), Inches(6.25), Inches(1.9),
          'Future work', [
              'Apply models in DCT / wavelet domain — images are',
              'sparser there, closer to the Gaussian-sparse model.',
              '',
              'Design image-aware thresholds for bounded,',
              'non-negative pixel structure.',
          ],
          left_accent=ACCENT, bg=RGBColor(0xE8, 0xF5, 0xFF),
          border_color=RGBColor(0xC0, 0xE0, 0xFF), body_size=12)


def slide_16_references(prs):
    slide = _add_slide(prs)
    _slide_bg(slide, BG)
    _top_bar(slide)
    _bottom_bar(slide)
    _add_title(slide, 'Code & References')
    _page_number(slide, 16)

    _text_box(slide, Inches(0.55), Inches(1.3), W - Inches(1.1), Inches(0.35),
              'Code — public GitHub repository (notebooks 01–04 + src modules):',
              size=16, bold=True, color=BGU_BLUE)
    _text_box(slide, Inches(0.85), Inches(1.72), W - Inches(1.7), Inches(0.35),
              'github.com/ThEpiCake/hyperlista-mbdl', size=16, color=ACCENT)

    _text_box(slide, Inches(0.55), Inches(2.35), W - Inches(1.1), Inches(0.35),
              'References:', size=16, bold=True, color=BGU_BLUE)

    refs = [
        'A. Beck and M. Teboulle, "A Fast Iterative Shrinkage-Thresholding Algorithm for Linear Inverse Problems," SIAM J. Imaging Sciences, 2009.',
        'K. Gregor and Y. LeCun, "Learning Fast Approximations of Sparse Coding," ICML 2010.',
        'X. Chen, J. Liu, Z. Wang, W. Yin, "Theoretical Linear Convergence of Unfolded ISTA and Its Practical Weights and Thresholds," NeurIPS 2018.',
        'J. Liu, X. Chen, Z. Wang, W. Yin, "ALISTA: Analytic Weights Are As Good As Learned Weights in LISTA," ICLR 2019.',
        'X. Chen, J. Liu, Z. Wang, W. Yin, "Hyperparameter Tuning is All You Need for LISTA," NeurIPS 2021.',
        'N. Shlezinger, J. Whang, Y. C. Eldar, A. G. Dimakis, "Model-Based Deep Learning," Proceedings of the IEEE, 2023.',
        'Course lecture notes 1–8, Model-Based Deep Learning (361.2.2320), Ben-Gurion University, 2026.',
    ]
    for i, ref in enumerate(refs):
        _text_box(slide, Inches(0.85), Inches(2.8) + i * Inches(0.52),
                  W - Inches(1.7), Inches(0.5), f'[{i+1}]  {ref}', size=14, color=TEXT)

    _text_box(slide, Inches(0.55), Inches(6.55), W - Inches(1.1), Inches(0.4),
              'Thank you!  Questions welcome.', size=18, bold=True, color=BGU_BLUE)


# ═══════════════════════════════════════════════════════════════════════════════
# Main
# ═══════════════════════════════════════════════════════════════════════════════

# ── Speaker notes (Hebrew narration, ~15 minutes total) ──────────────────────
SPEAKER_NOTES = [
    # 1 Title (~30s)
    "שלום לכולם, אנחנו איתי ברון ואיתי ויגמן. הפרויקט שלנו עוסק בשחזור אותות דלילים "
    "ובחישה דחוסה בגישה של למידה עמוקה מבוססת מודל. שאלת הליבה שלנו: כמה רחוק אפשר להגיע "
    "עם מבנה אלגוריתמי בלבד, ומתי חייבים גמישות נלמדת מהדאטה. נקודת המוצא היא המאמר "
    "HyperLISTA מ-NeurIPS 2021.",
    # 2 Roadmap (~40s)
    "מבנה ההרצאה: נתחיל ברקע — בעיית LASSO והאלגוריתמים הקלאסיים ISTA ו-FISTA. משם נעבור "
    "ל-Deep Unrolling: LISTA, ALISTA ו-HyperLISTA. החלק המרכזי הוא המחקר שלנו — אילו רכיבים "
    "של ISTA שווה ללמוד, והשוואת שיטות אימון לרשתות unfolded. נסיים בניסויי תמונות אמיתיות "
    "על MNIST ו-FashionMNIST בחישה דחוסה. הסטאפ הסינתטי: 250 מדידות, ממד 500, דלילות 50.",
    # 3 Motivation (~60s)
    "בעיית השחזור: יש לנו וקטור דליל x כוכב בממד n, ואנחנו רואים רק m מדידות לינאריות, "
    "m קטן מ-n, עם רעש אפשרי. הבעיה לא מוגדרת היטב — יש אינסוף פתרונות — אבל הנחת הדלילות "
    "הופכת אותה לפתירה. זה מופיע בכל מקום: MRI, מכ\"ם, חישה דחוסה. הפורמליזציה הסטנדרטית "
    "היא LASSO: איזון בין נאמנות למדידות לבין נורמת L1 שמעודדת דלילות. המדדים שלנו: NMSE "
    "בדציבלים לאותות סינתטיים, ובנוסף PSNR ו-SSIM לתמונות.",
    # 4 Baselines (~60s)
    "ISTA הוא gradient descent פרוקסימלי על LASSO: צעד גרדיאנט על איבר הנאמנות, ואז "
    "soft-thresholding שמאפס רכיבים קטנים. FISTA מוסיף מומנטום של נסטרוב ומשפר את קצב "
    "ההתכנסות מסדר 1 חלקי k לסדר 1 חלקי k בריבוע. אבל שימו לב לנקודה המרכזית: אחרי 16 "
    "איטרציות בלבד ISTA מגיע רק למינוס 5.3 דציבל ו-FISTA למינוס 11. כדי להגיע למינוס 20 "
    "צריך מאות איטרציות. השאלה: איך גורמים ל-16 צעדים באמת לעבוד?",
    # 5 Unrolling (~70s)
    "הרעיון של Gregor ו-LeCun מ-2010: לפרוש את K האיטרציות של ISTA לרשת בעומק K, ולהפוך "
    "את הקבועים האלגוריתמיים לפרמטרים נלמדים. משמאל — ISTA קלאסי עם צעד ו-threshold קבועים. "
    "מימין — LISTA: שתי מטריצות משקל ו-threshold נפרד לכל שכבה, הכל נלמד ב-backpropagation "
    "מזוגות אימון. זה העיקרון של למידה מבוססת מודל: שומרים על המבנה — צעד גרדיאנט ואז "
    "soft-threshold — ולומדים רק את מה שמשפר שחזור בעומק סופי. בסטאפ שלנו זה כ-6 מיליון פרמטרים.",
    # 6 Ladder (~90s)
    "עכשיו סולם ה-MBDL — פחות פרמטרים, יותר מבנה. LISTA לומד הכל: 6 מיליון פרמטרים ומגיע "
    "למינוס 20 דציבל. ALISTA מראה שאפשר לחשב את מטריצת המשקל אנליטית מ-A — פתרון בעיית "
    "מזעור קוהרנטיות — ולהשאיר ללמידה רק צעדים ו-thresholds: 32 פרמטרים, מינוס 30 דציבל. "
    "HyperLISTA לוקח את זה לקצה: שלושה סקלרים בלבד. ה-threshold, המומנטום ובחירת התומך "
    "מחושבים אדפטיבית מהשארית בכל שכבה, לפי הנוסחאות שמוצגות. שימו לב — אין בכלל "
    "backpropagation: שלושת הסקלרים נמצאים ב-grid search. והתוצאה: מינוס 61 דציבל, הטובה "
    "ביותר. יותר מבנה, פחות פרמטרים, ביצועים טובים יותר — זה עקרון ה-MBDL בפעולה.",
    # 7 Design space (~60s)
    "זה מוביל לשאלת המחקר שלנו: בין 6 מיליון לשלושה פרמטרים — מה המינימום שבאמת צריך ללמוד? "
    "עדכון ISTA המוכלל חושף שתי סדרות סקלרים חופשיות: גודל הצעד גמא והסף תטא. בנינו שלושה "
    "מודלים: ThresholdISTA שלומד רק את הספים, StepISTA שלומד רק את הצעדים כשהסף צמוד אליהם, "
    "ו-StepThresholdISTA שלומד את שניהם בנפרד — 32 סקלרים בסך הכל. כולם מאותחלים לערכי ISTA "
    "המדויקים, כך שלפני האימון הם זהים ל-ISTA.",
    # 8 Scalar results (~70s)
    "התוצאות מפתיעות בעוצמתן. ThresholdISTA — לימוד ספים בלבד — כמעט לא עוזר: מינוס 7 דציבל. "
    "אבל StepISTA — לימוד לוחות הצעדים בלבד, 16 סקלרים — קופץ למינוס 20 דציבל ומשתווה ל-LISTA "
    "המלא עם 6 מיליון פרמטרים. הפרדת הסף מהצעד מוסיפה עוד 3 דציבל בחינם. המסקנה: לוח הצעדים "
    "הוא הרכיב הקריטי ללמידה בעומק סופי, ולמידת מטריצות מלאות היא בעיקרה מיותרת כשמטריצת "
    "המדידה קבועה.",
    # 9 Training methods (~70s)
    "בהתאם להמלצת המרצה, השווינו גם את שיטות האימון מההרצאות לאותה ארכיטקטורה בדיוק. "
    "L1 — אימון end-to-end עם loss על השכבה האחרונה. L2 — deep supervision עם loss משוקלל "
    "על כל שכבת ביניים. L3 — אימון חמדני שכבה-שכבה כשהקודמות מוקפאות. בנוסף בדקנו קשירת "
    "משקולות — אותה שכבה חוזרת K פעמים, כמו RNN. הממצא המעניין: LISTA קשור עם פי 16 פחות "
    "פרמטרים מגיע כמעט לאותה תוצאה — קשירת משקולות היא רגולריזציה מובנית.",
    # 10 Sparse results (~60s)
    "כאן התמונה המלאה על הבעיה הסינתטית. בגרף רואים NMSE כפונקציה של עומק. שיטות האימון "
    "משנות 3-4 דציבל בתוך משפחת LISTA — L2 הכי טוב, L3 נתקע ברמה מסוימת כי לשכבות מאוחרות "
    "אין תמריץ להשתפר. אבל ההיררכיה הגדולה נשארת: ALISTA מנצח את כל גרסאות LISTA עם 32 "
    "פרמטרים, ו-HyperLISTA מנצח את כולם עם 3. כשהנחות המודל נכונות — מבנה מנצח קיבולת.",
    # 11 Design map (~50s)
    "מפת ה-design space מסכמת את החלק הסינתטי: ציר X — מספר פרמטרים בסקלה לוגריתמית, ציר Y — "
    "שגיאת שחזור. הפינה השמאלית-תחתונה היא האידיאל. StepThresholdISTA שלנו — 32 סקלרים — "
    "יושב באזור מצוין ומנצח את LISTA עם פי 187 אלף פחות פרמטרים. HyperLISTA שובר את הסקלה "
    "לגמרי. המסר: בבעיות עם מבנה ידוע, תכנון חכם של מה לומדים שווה יותר מקיבולת.",
    # 12 Images (~70s)
    "עכשיו המבחן האמיתי — תמונות אמיתיות. עברנו לחישה דחוסה בפיקסלים על FashionMNIST: "
    "התמונה שטוחה לווקטור בממד 784, ומודדים רק רבע מהממד. אין צורך ב-DCT כי לתמונות רקע "
    "שחור והן דלילות בפיקסלים. והנה ההיפוך הדרמטי: כל השיטות המובנות קורסות — ALISTA "
    "ו-HyperLISTA גרועים אף מ-ISTA — בעוד LISTA-L2 מגיע ל-SSIM של 0.82. הסיבה: ההנחה של "
    "דלילות גאוסית i.i.d. פשוט לא נכונה לפיקסלים חסומים ומתואמים מרחבית. פריור שגוי גרוע "
    "מאין-פריור.",
    # 13 Robustness (~60s)
    "בדקנו את זה שיטתית על שני דאטהסטים ושני יחסי מדידה. מפת החום מציגה SSIM. "
    "על MNIST עם חצי מהמדידות — כולם סבירים, כי MNIST דליל באמת: 78 אחוז אפסים. ככל "
    "שמקשים — פחות מדידות או דאטה מרקמי יותר — השיטות המובנות מתדרדרות בחדות, בעוד "
    "LISTA-L2 נשאר מעל 0.82 בכל התרחישים. הפריור הדליל מתפרק בהדרגה כשהדאטה מתרחק מההנחות.",
    # 14 Stress tests (~50s)
    "ההשוואה הוויזואלית ממחישה את המספרים. משמאל MNIST, מימין FashionMNIST, שניהם ברבע "
    "מהמדידות. שורת LISTA-L2 כמעט זהה למקור בשני המקרים. ALISTA ו-HyperLISTA מאפסים סיגנל "
    "אמיתי — הסף שמכויל לרעש גאוסי דליל מוחק את התמונה עצמה. רואים כאן ויזואלית מה זה "
    "אומר כשהפריור שגוי.",
    # 15 Conclusions (~60s)
    "לסיכום, ארבעה לקחים: אחד — כשההנחות מתקיימות, מבנה מנצח: 3 סקלרים מנצחים 6 מיליון "
    "פרמטרים ב-41 דציבל. שתיים — לוח הצעדים הוא הרכיב הכי חשוב ללמידה. שלוש — קשירת "
    "משקולות היא רגולריזציה מובנית. ארבע — פריור שגוי מזיק יותר מאין-פריור: על תמונות "
    "אמיתיות LISTA הגמיש מנצח. בהמשך נרצה לבדוק את המודלים בדומיין DCT או wavelets, שם "
    "התמונות קרובות יותר להנחת הדלילות הגאוסית.",
    # 16 References (~20s)
    "כל הקוד, ארבעת המחברות והמודולים זמינים בריפו ציבורי בגיטהאב. אלו המקורות המרכזיים "
    "שעליהם התבססנו. תודה רבה על ההקשבה — נשמח לשאלות.",
]


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    print('Building slides...')
    builders = [
        slide_01_title,
        slide_02_roadmap,
        slide_03_motivation,
        slide_04_baselines,
        slide_05_unrolling,
        slide_06_mbdl_ladder,
        slide_07_design,
        slide_08_scalar_results,
        slide_09_training,
        slide_10_sparse_results,
        slide_11_design_space,
        slide_12_images,
        slide_13_robustness,
        slide_14_stress,
        slide_14_conclusions,
        slide_16_references,
    ]
    for i, fn in enumerate(builders, 1):
        print(f'  Slide {i:2d}: {fn.__name__}')
        fn(prs)

    # Attach speaker notes
    for slide, note in zip(prs.slides, SPEAKER_NOTES):
        slide.notes_slide.notes_text_frame.text = note

    out = os.path.join(os.path.dirname(__file__), 'MBDL_Seminar_Final.pptx')
    out = os.path.normpath(out)
    prs.save(out)
    print(f'\nSaved: {out}')


if __name__ == '__main__':
    build()
